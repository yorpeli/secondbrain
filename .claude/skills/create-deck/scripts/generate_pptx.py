#!/usr/bin/env python3
"""
PowerPoint Generation Script for /create-deck skill.
Generates branded PPTX files with elegant glass-effect design system.
"""

import argparse
import json
import sys
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from lxml import etree
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


# Global color constants for glass effects
CHARCOAL_BLACK = hex_to_rgb("#1E1E28")
PURE_WHITE = hex_to_rgb("#FFFFFF")
LIGHT_GRAY = hex_to_rgb("#E5E5E5")
MUTED_TEXT = hex_to_rgb("#666666")
CONTAINER_FILL = hex_to_rgb("#FAFAFA")
SHADOW_COLOR = hex_to_rgb("#E0E0E0")
ELECTRIC_BLUE = hex_to_rgb("#0033FF")
MIDNIGHT_BLUE = hex_to_rgb("#002373")
NEON_PURPLE = hex_to_rgb("#977DFF")

# Glass effect colors (from reference design)
GLASS_GRADIENT_END = "#F2FAFF"  # Very light blue for gradient edge
GLASS_BORDER_COLOR = "#D9D9D9"  # Light gray border (85% white)


def apply_glass_effect(shape, gradient_angle=315, border_width_pt=2.0,
                       shadow_blur_pt=100, shadow_dist_pt=20, shadow_alpha_pct=12):
    """
    Apply elegant glass effect to a shape using oxml.

    Creates:
    - Linear gradient fill (white center → light blue edge)
    - Light gray border
    - Subtle outer shadow with blur
    """
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))

    if spPr is None:
        return

    # Remove existing fill
    for fill in spPr.findall(qn('a:solidFill')):
        spPr.remove(fill)
    for fill in spPr.findall(qn('a:gradFill')):
        spPr.remove(fill)
    for fill in spPr.findall(qn('a:noFill')):
        spPr.remove(fill)

    # Create gradient fill XML
    # Angle in EMUs: 315° * 60000 = 18900000
    angle_emu = int(gradient_angle * 60000)

    gradFill_xml = f'''
    <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:gsLst>
            <a:gs pos="46000">
                <a:srgbClr val="FFFFFF"/>
            </a:gs>
            <a:gs pos="100000">
                <a:srgbClr val="F2FAFF"/>
            </a:gs>
        </a:gsLst>
        <a:lin ang="{angle_emu}" scaled="1"/>
    </a:gradFill>
    '''
    gradFill = etree.fromstring(gradFill_xml)

    # Insert gradient fill after xfrm and prstGeom
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.addnext(gradFill)
    else:
        spPr.insert(0, gradFill)

    # Update line properties
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(ln)

    # Border width in EMUs: pt * 12700
    border_emu = int(border_width_pt * 12700)

    ln_xml = f'''
    <a:ln w="{border_emu}" cap="flat" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:solidFill>
            <a:srgbClr val="D9D9D9"/>
        </a:solidFill>
        <a:prstDash val="solid"/>
        <a:miter lim="800000"/>
    </a:ln>
    '''
    ln = etree.fromstring(ln_xml)
    spPr.append(ln)

    # Add outer shadow effect
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is not None:
        spPr.remove(effectLst)

    # Shadow values in EMUs
    blur_emu = int(shadow_blur_pt * 12700)
    dist_emu = int(shadow_dist_pt * 12700)
    alpha_val = int(shadow_alpha_pct * 1000)  # 12% = 12000

    effectLst_xml = f'''
    <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="{blur_emu}" dist="{dist_emu}" dir="5400000" sx="90000" sy="90000" algn="ctr" rotWithShape="0">
            <a:srgbClr val="002373">
                <a:alpha val="{alpha_val}"/>
            </a:srgbClr>
        </a:outerShdw>
    </a:effectLst>
    '''
    effectLst = etree.fromstring(effectLst_xml)
    spPr.append(effectLst)

    # Update corner radius (adjustment value ~9.4% = 9356/100000)
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is not None:
            for gd in avLst.findall(qn('a:gd')):
                if gd.get('name') == 'adj':
                    gd.set('fmla', 'val 9356')


class DeckGenerator:
    """Generate branded PowerPoint presentations with glass-effect design."""

    def __init__(self, brand_config: dict):
        self.brand = brand_config
        self.prs = Presentation()

        # Set slide dimensions (widescreen 16:9)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # Cache colors
        self.colors = {
            'background': hex_to_rgb(self.brand['slide_defaults']['background_color']),
            'title': hex_to_rgb(self.brand['slide_defaults']['title_color']),
            'body': hex_to_rgb(self.brand['slide_defaults']['body_color']),
            'accent': hex_to_rgb(self.brand['slide_defaults']['accent_color']),
            'table_header_bg': hex_to_rgb(self.brand['slide_defaults']['table_header_bg']),
            'table_header_text': hex_to_rgb(self.brand['slide_defaults']['table_header_text']),
            'table_row_alt': hex_to_rgb(self.brand['slide_defaults']['table_row_alt_bg']),
        }

        # Font settings
        self.font_family = self.brand['typography'].get('fallback', 'Arial')
        self.font_sizes = self.brand['typography']['sizes']

        # Glass component settings
        self.glass_config = self.brand.get('glass_components', {})
        self.badge_styles = self.brand.get('badge_styles', {})
        self.slide_themes = self.brand.get('slide_themes', {})

    def _add_text(self, slide, text: str, left, top, width=None, height=None,
                  font_size: int = 16, color=None, bold: bool = False,
                  alignment=PP_ALIGN.LEFT, uppercase: bool = False):
        """Helper to add formatted text to a slide."""
        width = width or Inches(10)
        height = height or Inches(0.5)
        color = color or self.colors['body']

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        para = text_frame.paragraphs[0]
        para.text = text.upper() if uppercase else text
        para.font.size = Pt(font_size)
        para.font.name = self.font_family
        para.font.color.rgb = color
        para.font.bold = bold
        para.alignment = alignment

        return text_box

    def add_glass_card(self, slide, left, top, width, height, with_shadow: bool = True):
        """
        Add an elegant glass-effect container.

        Creates a rounded rectangle with:
        - Linear gradient fill (white center → light blue #F2FAFF edge at 315°)
        - Light gray border (2pt)
        - Subtle outer shadow with blur (if enabled)
        - ~9.4% corner radius
        """
        card_config = self.glass_config.get('card', {})

        # Create the glass card shape
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )

        # Set initial corner radius (will be refined by apply_glass_effect)
        if hasattr(card, 'adjustments') and len(card.adjustments) > 0:
            card.adjustments[0] = 0.094  # ~9.4%

        # Apply the glass effect (gradient fill, border, shadow)
        apply_glass_effect(
            card,
            gradient_angle=315,
            border_width_pt=card_config.get('border_weight_pt', 2.0),
            shadow_blur_pt=100 if with_shadow else 0,
            shadow_dist_pt=20 if with_shadow else 0,
            shadow_alpha_pct=12 if with_shadow else 0
        )

        return card

    def add_border_card(self, slide, left, top, width, height):
        """Add a border-only card (no fill)."""
        card_config = self.glass_config.get('border_only_card', {})
        border_weight = Pt(card_config.get('border_weight_pt', 1.0))

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        card.fill.background()
        card.line.color.rgb = LIGHT_GRAY
        card.line.width = border_weight
        if hasattr(card, 'adjustments') and len(card.adjustments) > 0:
            card.adjustments[0] = 0.05

        return card

    def add_badge(self, slide, text: str, left, top, style: str = "primary"):
        """Add a status badge (pill shape)."""
        styles = self.badge_styles or {
            "primary": {"fill": "#002373", "text_color": "#FFFFFF"},
            "secondary": {"fill": "#977DFF", "text_color": "#FFFFFF"},
            "outline": {"fill": "#FAFAFA", "text_color": "#1E1E28", "border_color": "#1E1E28"},
            "neutral": {"fill": "#F5F5F5", "text_color": "#666666"},
        }

        style_config = styles.get(style, styles.get("primary", {}))
        badge_font_size = self.font_sizes.get('badge', 10)

        # Calculate badge size based on text
        width = Inches(max(0.6, len(text) * 0.08 + 0.3))
        height = Inches(0.25)

        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )

        # Apply fill
        fill_color = style_config.get('fill')
        if fill_color:
            badge.fill.solid()
            badge.fill.fore_color.rgb = hex_to_rgb(fill_color)
        else:
            badge.fill.background()

        # Apply border if specified
        border_color = style_config.get('border_color')
        if border_color:
            badge.line.color.rgb = hex_to_rgb(border_color)
            badge.line.width = Pt(0.75)
        else:
            badge.line.fill.background()

        # Rounded corners (pill shape)
        if hasattr(badge, 'adjustments') and len(badge.adjustments) > 0:
            badge.adjustments[0] = 0.5

        # Add text to badge
        text_color = hex_to_rgb(style_config.get('text_color', '#FFFFFF'))
        badge.text_frame.paragraphs[0].text = text.upper()
        badge.text_frame.paragraphs[0].font.size = Pt(badge_font_size)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = text_color
        badge.text_frame.paragraphs[0].font.name = self.font_family
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.word_wrap = False

        return badge

    def add_section_header(self, slide, text: str, left, top, accent_color=None):
        """Add elegant section header with short underline."""
        color = accent_color or self.colors['accent']
        section_font_size = self.font_sizes.get('section_header', 12)

        # Uppercase text
        self._add_text(
            slide, text, left, top,
            width=Inches(4), height=Inches(0.35),
            font_size=section_font_size, color=color,
            bold=True, uppercase=True
        )

        # Short accent underline (0.6")
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top + Inches(0.3), Inches(0.6), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

        return line

    def add_accent_line(self, slide, left, top, width: float = 0.8):
        """Add a short accent line."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, Inches(width), Pt(3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent']
        line.line.fill.background()
        return line

    def add_dark_title_slide(self, title: str, subtitle: str = None, date: str = None):
        """Add elegant dark title slide with impact."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Dark background
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = CHARCOAL_BLACK

        # Short accent line above title
        self.add_accent_line(slide, Inches(0.8), Inches(2.8), width=0.8)

        # White title
        self._add_text(
            slide, title, Inches(0.8), Inches(3.0),
            width=Inches(11.73), height=Inches(1.2),
            font_size=self.font_sizes.get('title', 48),
            color=PURE_WHITE, bold=True
        )

        # Gray subtitle
        if subtitle:
            self._add_text(
                slide, subtitle, Inches(0.8), Inches(3.9),
                width=Inches(11.73), height=Inches(0.6),
                font_size=24, color=LIGHT_GRAY
            )

        # Date in muted text
        if date:
            self._add_text(
                slide, date, Inches(0.8), Inches(4.6),
                width=Inches(4), height=Inches(0.4),
                font_size=14, color=MUTED_TEXT
            )

        return slide

    def add_title_slide(self, title: str, subtitle: str = None, dark: bool = False):
        """Add a title slide (light or dark variant)."""
        if dark:
            return self.add_dark_title_slide(title, subtitle)

        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(2.5), Inches(11.83), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(self.font_sizes.get('title', 48))
        title_para.font.name = self.font_family
        title_para.font.color.rgb = self.colors['title']
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(4.2), Inches(11.83), Inches(0.75)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(self.font_sizes.get('body', 16))
            subtitle_para.font.name = self.font_family
            subtitle_para.font.color.rgb = self.colors['body']
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Add accent line (0.8" centered)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.27), Inches(4.0), Inches(0.8), Pt(3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent']
        line.line.fill.background()

        return slide

    def add_content_slide(self, title: str, bullets: list, highlight_first: bool = False,
                          subtitle: str = None, in_glass_card: bool = False):
        """Add a content slide with bullet points."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add slide title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(self.font_sizes.get('slide_title', 32))
        title_para.font.name = self.font_family
        title_para.font.color.rgb = self.colors['title']
        title_para.font.bold = True

        # Add muted subtitle if provided
        if subtitle:
            self._add_text(
                slide, subtitle, Inches(0.75), Inches(1.1),
                width=Inches(11.83), height=Inches(0.4),
                font_size=14, color=MUTED_TEXT
            )
            content_top = Inches(1.6)
        else:
            content_top = Inches(1.3)

        # Add accent underline (shorter, 0.8")
        self.add_accent_line(slide, Inches(0.75), content_top - Inches(0.25))

        # Add glass card if requested
        if in_glass_card:
            card_top = content_top + Inches(0.1)
            card_height = Inches(5.0)
            self.add_glass_card(
                slide, Inches(0.6), card_top,
                Inches(12.1), card_height
            )
            content_left = Inches(0.85)
            content_top = card_top + Inches(0.25)
        else:
            content_left = Inches(0.75)
            content_top = content_top + Inches(0.15)

        # Add bullet content
        content_box = slide.shapes.add_textbox(
            content_left, content_top, Inches(11.5), Inches(5.0)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()

            para.text = f"• {bullet}"
            para.font.size = Pt(self.font_sizes.get('body', 16))
            para.font.name = self.font_family
            para.space_after = Pt(12)

            # Highlight first bullet if requested
            if i == 0 and highlight_first:
                para.font.color.rgb = self.colors['accent']
                para.font.bold = True
            else:
                para.font.color.rgb = self.colors['body']

        return slide

    def add_content_cards_slide(self, title: str, cards: list, subtitle: str = None):
        """Add a slide with multiple glass card containers for grouped content."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add slide title
        self._add_text(
            slide, title, Inches(0.75), Inches(0.5),
            width=Inches(11.83), height=Inches(1.0),
            font_size=self.font_sizes.get('slide_title', 32),
            color=self.colors['title'], bold=True
        )

        # Add muted subtitle if provided
        if subtitle:
            self._add_text(
                slide, subtitle, Inches(0.75), Inches(1.1),
                width=Inches(11.83), height=Inches(0.4),
                font_size=14, color=MUTED_TEXT
            )
            self.add_accent_line(slide, Inches(0.75), Inches(1.5))
            cards_top = Inches(1.8)
        else:
            self.add_accent_line(slide, Inches(0.75), Inches(1.3))
            cards_top = Inches(1.6)

        # Calculate card layout
        num_cards = len(cards)
        if num_cards == 0:
            return slide

        total_width = Inches(12.0)
        gap = Inches(0.3)
        card_width = (total_width - gap * (num_cards - 1)) / num_cards if num_cards > 1 else total_width
        card_height = Inches(5.2)

        # Add each card
        for idx, card_data in enumerate(cards):
            card_left = Inches(0.65) + idx * (card_width + gap)

            # Add glass card background
            self.add_glass_card(slide, card_left, cards_top, card_width, card_height)

            # Card section header
            header_text = card_data.get('header', f'Section {idx + 1}')
            self.add_section_header(
                slide, header_text,
                card_left + Inches(0.2), cards_top + Inches(0.2)
            )

            # Card content
            content = card_data.get('content', [])
            content_box = slide.shapes.add_textbox(
                card_left + Inches(0.2), cards_top + Inches(0.7),
                card_width - Inches(0.4), card_height - Inches(0.9)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, bullet in enumerate(content):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()

                para.text = f"• {bullet}"
                para.font.size = Pt(self.font_sizes.get('body', 16))
                para.font.name = self.font_family
                para.font.color.rgb = self.colors['body']
                para.space_after = Pt(8)

        return slide

    def add_metrics_slide(self, title: str, metrics: list, subtitle: str = None):
        """Add a slide with metric cards (3-4 per row)."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add slide title
        self._add_text(
            slide, title, Inches(0.75), Inches(0.5),
            width=Inches(11.83), height=Inches(1.0),
            font_size=self.font_sizes.get('slide_title', 32),
            color=self.colors['title'], bold=True
        )

        # Add muted subtitle if provided
        if subtitle:
            self._add_text(
                slide, subtitle, Inches(0.75), Inches(1.1),
                width=Inches(11.83), height=Inches(0.4),
                font_size=14, color=MUTED_TEXT
            )
            self.add_accent_line(slide, Inches(0.75), Inches(1.5))
            metrics_top = Inches(1.9)
        else:
            self.add_accent_line(slide, Inches(0.75), Inches(1.3))
            metrics_top = Inches(1.7)

        # Calculate metric card layout
        num_metrics = min(len(metrics), 4)  # Max 4 per row
        if num_metrics == 0:
            return slide

        total_width = Inches(12.0)
        gap = Inches(0.4)
        card_width = (total_width - gap * (num_metrics - 1)) / num_metrics
        card_height = Inches(2.5)

        for idx, metric in enumerate(metrics[:4]):
            card_left = Inches(0.65) + idx * (card_width + gap)

            # Add glass card
            self.add_glass_card(slide, card_left, metrics_top, card_width, card_height)

            # Metric value (large, bold)
            value = metric.get('value', '0')
            self._add_text(
                slide, str(value), card_left + Inches(0.3), metrics_top + Inches(0.4),
                width=card_width - Inches(0.6), height=Inches(1.2),
                font_size=40, color=self.colors['accent'], bold=True,
                alignment=PP_ALIGN.CENTER
            )

            # Metric label
            label = metric.get('label', 'Metric')
            self._add_text(
                slide, label, card_left + Inches(0.3), metrics_top + Inches(1.5),
                width=card_width - Inches(0.6), height=Inches(0.6),
                font_size=self.font_sizes.get('body', 16),
                color=MUTED_TEXT, alignment=PP_ALIGN.CENTER
            )

            # Optional badge
            badge_text = metric.get('badge')
            if badge_text:
                badge_style = metric.get('badge_style', 'neutral')
                badge_left = card_left + (card_width - Inches(0.8)) / 2
                self.add_badge(slide, badge_text, badge_left, metrics_top + Inches(2.0), badge_style)

        return slide

    def add_image_slide(self, title: str, image_path: str, subtitle: str = None,
                        caption: str = None, in_glass_card: bool = False):
        """Add a slide with an embedded image, properly scaled."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add slide title
        self._add_text(
            slide, title, Inches(0.75), Inches(0.5),
            width=Inches(11.83), height=Inches(1.0),
            font_size=self.font_sizes.get('slide_title', 32),
            color=self.colors['title'], bold=True
        )

        # Add subtitle if provided
        if subtitle:
            self._add_text(
                slide, subtitle, Inches(0.75), Inches(1.1),
                width=Inches(11.83), height=Inches(0.4),
                font_size=14, color=MUTED_TEXT
            )
            self.add_accent_line(slide, Inches(0.75), Inches(1.5))
            img_area_top = Inches(1.8)
        else:
            self.add_accent_line(slide, Inches(0.75), Inches(1.3))
            img_area_top = Inches(1.6)

        # Calculate available area for image
        img_area_bottom = Inches(7.0) if not caption else Inches(6.5)
        img_area_height = img_area_bottom - img_area_top
        img_area_width = Inches(11.83)

        # Add glass card behind image if requested
        if in_glass_card:
            card_margin = Inches(0.15)
            self.add_glass_card(
                slide, Inches(0.6), img_area_top - card_margin,
                Inches(12.1), img_area_height + card_margin * 2
            )

        # Load image to get aspect ratio
        try:
            from PIL import Image as PILImage
            with PILImage.open(image_path) as img:
                img_w, img_h = img.size
        except Exception:
            # Fallback: assume 16:9 aspect ratio
            img_w, img_h = 1600, 900

        aspect = img_w / img_h
        available_aspect = img_area_width / img_area_height

        # Scale to fit within available area while preserving aspect ratio
        if aspect > available_aspect:
            # Image is wider — constrain by width
            width = img_area_width
            height = width / aspect
        else:
            # Image is taller — constrain by height
            height = img_area_height
            width = height * aspect

        # Center the image
        left = Inches(0.75) + (img_area_width - width) / 2
        top = img_area_top + (img_area_height - height) / 2

        slide.shapes.add_picture(image_path, int(left), int(top), int(width), int(height))

        # Add caption if provided
        if caption:
            self._add_text(
                slide, caption, Inches(0.75), Inches(6.6),
                width=Inches(11.83), height=Inches(0.4),
                font_size=self.font_sizes.get('caption', 12),
                color=MUTED_TEXT, alignment=PP_ALIGN.CENTER
            )

        return slide

    def add_closing_slide(self, title: str = "Questions?", subtitle: str = None):
        """Add a dark closing slide (Q&A style)."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Dark background
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = CHARCOAL_BLACK

        # Accent line
        self.add_accent_line(slide, Inches(5.87), Inches(3.0), width=1.6)

        # Title (centered)
        self._add_text(
            slide, title, Inches(0.75), Inches(3.3),
            width=Inches(11.83), height=Inches(1.0),
            font_size=self.font_sizes.get('title', 48),
            color=PURE_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER
        )

        # Subtitle
        if subtitle:
            self._add_text(
                slide, subtitle, Inches(0.75), Inches(4.3),
                width=Inches(11.83), height=Inches(0.6),
                font_size=18, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        return slide

    def add_dark_grid_table(self, title: str, headers: list, rows: list,
                            insights: list = None, subtitle: str = None):
        """
        Add a dark-background comparison grid built from text boxes.

        Mimics the MFA_Competitive_Analysis slide 7 style:
        - Dark background, accent line at top
        - Column headers: first in purple, rest in white
        - Colored symbols: checkmark (green), cross (red), tilde (orange)
        - Text values colored by sentiment
        - Optional insights section at bottom
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Dark background
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = CHARCOAL_BLACK

        # Accent line at top
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(0.6), Inches(0.6), Pt(2)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.colors['accent']
        accent_line.line.fill.background()

        # Title
        self._add_text(
            slide, title, Inches(0.8), Inches(0.75),
            width=Inches(10), height=Inches(0.8),
            font_size=self.font_sizes.get('slide_title', 32),
            color=PURE_WHITE, bold=True
        )

        # Subtitle if provided
        header_top = Inches(1.6)
        if subtitle:
            self._add_text(
                slide, subtitle, Inches(0.8), Inches(1.3),
                width=Inches(10), height=Inches(0.35),
                font_size=14, color=MUTED_TEXT
            )
            header_top = Inches(1.75)

        # Layout calculations
        num_cols = len(headers)
        left_margin = Inches(0.8)
        first_col_width = Inches(3.2)
        remaining_width = Inches(11.5) - first_col_width
        data_col_width = remaining_width / max(num_cols - 1, 1)
        row_height = Inches(0.5)

        # Column headers
        for col_idx, header_text in enumerate(headers):
            if col_idx == 0:
                col_left = left_margin
                col_w = first_col_width
                h_color = NEON_PURPLE
            else:
                col_left = left_margin + first_col_width + data_col_width * (col_idx - 1)
                col_w = data_col_width
                h_color = PURE_WHITE

            self._add_text(
                slide, header_text, col_left, header_top,
                width=col_w, height=Inches(0.4),
                font_size=14, color=h_color, bold=True,
                alignment=PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            )

        # Separator line below headers
        sep_top = header_top + Inches(0.45)
        sep_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left_margin, sep_top, Inches(11.5), Pt(1.5)
        )
        sep_line.fill.solid()
        sep_line.fill.fore_color.rgb = MUTED_TEXT
        sep_line.line.fill.background()

        # Data rows
        row_top = sep_top + Inches(0.2)

        # Color mapping for cell values
        GREEN = hex_to_rgb("#228B22")
        RED = hex_to_rgb("#DC143C")
        ORANGE = hex_to_rgb("#FF8C00")

        def get_cell_color_and_symbol(value: str):
            """Determine display text and color based on cell value."""
            val_lower = value.lower().strip()
            if val_lower.startswith("yes"):
                return value, GREEN
            elif val_lower in ("no", "no."):
                return "\u2717", RED
            elif val_lower.startswith("no "):
                # e.g., "No — IDE-locked"
                return value, RED
            elif val_lower in ("partial", "limited", "chat only"):
                return "~  " + value, ORANGE
            elif val_lower.startswith("limited"):
                return "~  " + value, ORANGE
            elif val_lower.startswith("chat"):
                return "~  " + value, ORANGE
            elif val_lower.startswith("partial"):
                return "~  " + value, ORANGE
            else:
                return value, LIGHT_GRAY

        for row_idx, row_data in enumerate(rows):
            y = row_top + row_height * row_idx

            for col_idx, cell_value in enumerate(row_data):
                if col_idx == 0:
                    # Row label
                    self._add_text(
                        slide, cell_value, left_margin, y,
                        width=first_col_width, height=row_height,
                        font_size=12, color=PURE_WHITE
                    )
                else:
                    col_left = left_margin + first_col_width + data_col_width * (col_idx - 1)
                    display_text, cell_color = get_cell_color_and_symbol(cell_value)
                    self._add_text(
                        slide, display_text, col_left, y,
                        width=data_col_width, height=row_height,
                        font_size=13, color=cell_color, bold=True,
                        alignment=PP_ALIGN.CENTER
                    )

        # Insights section
        if insights:
            insights_top = row_top + row_height * len(rows) + Inches(0.25)

            # Separator
            sep2 = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left_margin, insights_top, Inches(11.5), Pt(1.5)
            )
            sep2.fill.solid()
            sep2.fill.fore_color.rgb = MUTED_TEXT
            sep2.line.fill.background()

            # "Key Insights" header
            self._add_text(
                slide, "Key Insights", left_margin, insights_top + Inches(0.15),
                width=Inches(10), height=Inches(0.4),
                font_size=16, color=NEON_PURPLE, bold=True
            )

            # Insight bullets
            for idx, insight in enumerate(insights):
                self._add_text(
                    slide, f"\u2022 {insight}", left_margin, insights_top + Inches(0.5) + Inches(0.28) * idx,
                    width=Inches(11.5), height=Inches(0.3),
                    font_size=11, color=PURE_WHITE
                )

        return slide

    def add_table_slide(self, title: str, headers: list, rows: list, subtitle: str = None):
        """Add a slide with an elegantly styled comparison table."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add slide title
        self._add_text(
            slide, title, Inches(0.75), Inches(0.5),
            width=Inches(11.83), height=Inches(1.0),
            font_size=self.font_sizes.get('slide_title', 32),
            color=self.colors['title'], bold=True
        )

        # Add subtitle if provided
        if subtitle:
            self._add_text(
                slide, subtitle, Inches(0.75), Inches(1.1),
                width=Inches(11.83), height=Inches(0.4),
                font_size=14, color=MUTED_TEXT
            )
            table_top = Inches(1.7)
        else:
            table_top = Inches(1.5)

        # Calculate table dimensions
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header

        table_width = Inches(11.83)
        table_height = Inches(5.3) if not subtitle else Inches(5.0)
        col_width = table_width / num_cols

        # Add table
        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.75), table_top,
            table_width, table_height
        ).table

        # Style header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['table_header_bg']

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(self.font_sizes.get('body', 16))
            para.font.name = self.font_family
            para.font.color.rgb = self.colors['table_header_text']
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Style data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_value)

                # Alternate row colors (subtle glass effect)
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CONTAINER_FILL
                else:
                    cell.fill.background()

                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(self.font_sizes.get('caption', 12))
                para.font.name = self.font_family
                para.font.color.rgb = self.colors['body']
                para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        return slide

    def add_two_column_slide(self, title: str, left_title: str, left_bullets: list,
                             right_title: str, right_bullets: list, in_glass_cards: bool = True):
        """Add a two-column content slide with optional glass cards."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add slide title
        self._add_text(
            slide, title, Inches(0.75), Inches(0.5),
            width=Inches(11.83), height=Inches(1.0),
            font_size=self.font_sizes.get('slide_title', 32),
            color=self.colors['title'], bold=True
        )

        # Add accent line
        self.add_accent_line(slide, Inches(0.75), Inches(1.3))

        columns_top = Inches(1.6)
        col_width = Inches(5.8)
        col_height = Inches(5.2)

        # Add glass cards if enabled
        if in_glass_cards:
            self.add_glass_card(slide, Inches(0.6), columns_top, col_width, col_height)
            self.add_glass_card(slide, Inches(6.9), columns_top, col_width, col_height)
            content_offset = Inches(0.2)
        else:
            content_offset = Inches(0)

        # Left column header
        self.add_section_header(
            slide, left_title,
            Inches(0.75) + content_offset, columns_top + Inches(0.15)
        )

        # Left column content
        left_content_box = slide.shapes.add_textbox(
            Inches(0.75) + content_offset, columns_top + Inches(0.6),
            col_width - Inches(0.4), col_height - Inches(0.8)
        )
        lc_frame = left_content_box.text_frame
        lc_frame.word_wrap = True

        for i, bullet in enumerate(left_bullets):
            if i == 0:
                para = lc_frame.paragraphs[0]
            else:
                para = lc_frame.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(self.font_sizes.get('body', 16))
            para.font.name = self.font_family
            para.font.color.rgb = self.colors['body']
            para.space_after = Pt(8)

        # Right column header
        self.add_section_header(
            slide, right_title,
            Inches(7.05) + content_offset, columns_top + Inches(0.15)
        )

        # Right column content
        right_content_box = slide.shapes.add_textbox(
            Inches(7.05) + content_offset, columns_top + Inches(0.6),
            col_width - Inches(0.4), col_height - Inches(0.8)
        )
        rc_frame = right_content_box.text_frame
        rc_frame.word_wrap = True

        for i, bullet in enumerate(right_bullets):
            if i == 0:
                para = rc_frame.paragraphs[0]
            else:
                para = rc_frame.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(self.font_sizes.get('body', 16))
            para.font.name = self.font_family
            para.font.color.rgb = self.colors['body']
            para.space_after = Pt(8)

        return slide

    def generate_from_json(self, content: dict) -> int:
        """Generate deck from structured JSON content."""
        slides_created = 0

        for slide_data in content.get('slides', []):
            slide_type = slide_data.get('type', 'content')

            if slide_type == 'title':
                self.add_title_slide(
                    slide_data.get('title', 'Untitled'),
                    slide_data.get('subtitle'),
                    dark=slide_data.get('dark', False)
                )
                slides_created += 1

            elif slide_type == 'title_dark':
                date_str = slide_data.get('date', datetime.now().strftime('%B %Y'))
                self.add_dark_title_slide(
                    slide_data.get('title', 'Untitled'),
                    slide_data.get('subtitle'),
                    date_str
                )
                slides_created += 1

            elif slide_type == 'content':
                bullets = slide_data.get('content', [])
                if bullets:
                    self.add_content_slide(
                        slide_data.get('title', ''),
                        bullets,
                        slide_data.get('highlight_first', False),
                        slide_data.get('subtitle'),
                        slide_data.get('in_glass_card', False)
                    )
                    slides_created += 1

            elif slide_type == 'content_cards':
                cards = slide_data.get('cards', [])
                if cards:
                    self.add_content_cards_slide(
                        slide_data.get('title', ''),
                        cards,
                        slide_data.get('subtitle')
                    )
                    slides_created += 1

            elif slide_type == 'metrics':
                metrics = slide_data.get('metrics', [])
                if metrics:
                    self.add_metrics_slide(
                        slide_data.get('title', 'Key Metrics'),
                        metrics,
                        slide_data.get('subtitle')
                    )
                    slides_created += 1

            elif slide_type == 'table':
                table_data = slide_data.get('table_data', {})
                if table_data:
                    self.add_table_slide(
                        slide_data.get('title', 'Comparison'),
                        table_data.get('headers', []),
                        table_data.get('rows', []),
                        slide_data.get('subtitle')
                    )
                    slides_created += 1

            elif slide_type == 'table_dark':
                table_data = slide_data.get('table_data', {})
                if table_data:
                    self.add_dark_grid_table(
                        slide_data.get('title', 'Comparison'),
                        table_data.get('headers', []),
                        table_data.get('rows', []),
                        slide_data.get('insights'),
                        slide_data.get('subtitle')
                    )
                    slides_created += 1

            elif slide_type == 'two_column':
                self.add_two_column_slide(
                    slide_data.get('title', ''),
                    slide_data.get('left_title', 'Left'),
                    slide_data.get('left_content', []),
                    slide_data.get('right_title', 'Right'),
                    slide_data.get('right_content', []),
                    slide_data.get('in_glass_cards', True)
                )
                slides_created += 1

            elif slide_type == 'image':
                image_path = slide_data.get('image_path')
                if image_path and Path(image_path).exists():
                    self.add_image_slide(
                        slide_data.get('title', ''),
                        image_path,
                        slide_data.get('subtitle'),
                        slide_data.get('caption'),
                        slide_data.get('in_glass_card', False)
                    )
                    slides_created += 1
                else:
                    print(f"Warning: Image not found: {image_path}")

            elif slide_type == 'closing':
                self.add_closing_slide(
                    slide_data.get('title', 'Questions?'),
                    slide_data.get('subtitle')
                )
                slides_created += 1

        return slides_created

    def save(self, output_path: str):
        """Save the presentation to file."""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        return str(output_path)


def main():
    parser = argparse.ArgumentParser(description='Generate branded PowerPoint presentations')
    parser.add_argument('--input', required=True, help='Path to JSON input file')
    parser.add_argument('--output', required=True, help='Path for output PPTX file')
    parser.add_argument('--brand', help='Path to brand config JSON (optional)')

    args = parser.parse_args()

    # Load input JSON
    try:
        with open(args.input, 'r') as f:
            content = json.load(f)
    except Exception as e:
        print(f"Error loading input file: {e}")
        sys.exit(1)

    # Load brand config from input or separate file
    brand_config = content.get('brand')
    if args.brand:
        try:
            with open(args.brand, 'r') as f:
                brand_config = json.load(f)
        except Exception as e:
            print(f"Error loading brand config: {e}")
            sys.exit(1)

    if not brand_config:
        print("Error: No brand configuration provided")
        sys.exit(1)

    # Generate presentation
    generator = DeckGenerator(brand_config)
    slides_created = generator.generate_from_json(content)
    output_path = generator.save(args.output)

    print(f"Generated {slides_created} slides")
    print(f"Saved to: {output_path}")


if __name__ == '__main__':
    main()
